#!/usr/bin/env python3
"""
Generate MS2 Demo Presentation (PPTX).

Creates a professional presentation covering the FACIS FAP MS2
end-to-end demonstrator: Simulation -> ORCE -> Kafka -> Validation.
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# Brand colors
DARK_BLUE = RGBColor(0x0B, 0x1D, 0x3A)  # Primary background
ACCENT_BLUE = RGBColor(0x1A, 0x73, 0xE8)  # Accent / links
LIGHT_BLUE = RGBColor(0x42, 0x85, 0xF4)  # Secondary accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GREEN = RGBColor(0x34, 0xA8, 0x53)  # Success
ORANGE = RGBColor(0xFB, 0xBC, 0x04)  # Warning / highlight
RED_ACCENT = RGBColor(0xEA, 0x43, 0x35)  # Alert

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_dark_background(slide):
    """Add a dark blue background to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE


def add_light_background(slide):
    """Add a light background to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_shape(slide, left, top, width, height, color, corner_radius=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    color=WHITE,
    bold=False,
    alignment=PP_ALIGN.LEFT,
    font_name="Calibri",
):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(
    slide,
    left,
    top,
    width,
    height,
    lines,
    font_size=16,
    color=WHITE,
    line_spacing=1.3,
    font_name="Calibri",
):
    """Add multi-line text with consistent formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(font_size * 0.4)
    return txBox


def add_bullet_list(
    slide, left, top, width, height, items, font_size=16, color=WHITE, bullet_char="\u2022"
):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(font_size * 0.3)
    return txBox


def create_title_slide(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_dark_background(slide)

    # Accent line
    add_shape(slide, Inches(1), Inches(2.3), Inches(1.5), Pt(4), ACCENT_BLUE)

    add_text_box(
        slide,
        Inches(1),
        Inches(2.5),
        Inches(10),
        Inches(1.2),
        "FACIS FAP",
        font_size=20,
        color=ACCENT_BLUE,
        bold=True,
    )

    add_text_box(
        slide,
        Inches(1),
        Inches(3.0),
        Inches(11),
        Inches(1.5),
        "MS2 End-to-End Demonstrator",
        font_size=44,
        color=WHITE,
        bold=True,
    )

    add_text_box(
        slide,
        Inches(1),
        Inches(4.3),
        Inches(10),
        Inches(1),
        "Simulation  \u2192  ORCE  \u2192  Kafka  \u2192  Validation",
        font_size=24,
        color=LIGHT_BLUE,
        bold=False,
    )

    add_text_box(
        slide,
        Inches(1),
        Inches(5.5),
        Inches(10),
        Inches(0.8),
        "IoT & AI over Trusted Zones  |  ATLAS IoT Lab GmbH  |  February 2026",
        font_size=16,
        color=MEDIUM_GRAY,
    )


def create_agenda_slide(prs):
    """Slide 2: Agenda."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_text_box(
        slide,
        Inches(1),
        Inches(0.5),
        Inches(10),
        Inches(0.8),
        "Agenda",
        font_size=36,
        color=WHITE,
        bold=True,
    )
    add_shape(slide, Inches(1), Inches(1.2), Inches(1), Pt(3), ACCENT_BLUE)

    items = [
        "MS2 Goal & Scope",
        "Architecture: Simulation \u2192 ORCE \u2192 Kafka",
        "Smart Energy Use Case",
        "Smart City Use Case",
        "End-to-End Data Flow (Live Demo)",
        "Validation Results",
        "Next Steps",
    ]

    for i, item in enumerate(items):
        y = Inches(1.8) + Inches(i * 0.7)
        # Number circle
        num_shape = add_shape(slide, Inches(1.2), y, Inches(0.45), Inches(0.45), ACCENT_BLUE)
        num_shape.text_frame.paragraphs[0].text = str(i + 1)
        num_shape.text_frame.paragraphs[0].font.size = Pt(16)
        num_shape.text_frame.paragraphs[0].font.color.rgb = WHITE
        num_shape.text_frame.paragraphs[0].font.bold = True
        num_shape.text_frame.paragraphs[0].font.name = "Calibri"
        num_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        num_shape.text_frame.paragraphs[0].space_before = Pt(2)

        add_text_box(slide, Inches(2.0), y, Inches(8), Inches(0.5), item, font_size=20, color=WHITE)


def create_goal_slide(prs):
    """Slide 3: MS2 Goal & Scope."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_text_box(
        slide,
        Inches(1),
        Inches(0.5),
        Inches(10),
        Inches(0.8),
        "MS2 Goal & Scope",
        font_size=36,
        color=WHITE,
        bold=True,
    )
    add_shape(slide, Inches(1), Inches(1.2), Inches(1), Pt(3), ACCENT_BLUE)

    # Goal box
    goal_box = add_shape(
        slide, Inches(1), Inches(1.6), Inches(11), Inches(1.2), RGBColor(0x12, 0x2A, 0x4C)
    )
    add_text_box(
        slide,
        Inches(1.3),
        Inches(1.7),
        Inches(10),
        Inches(0.4),
        "GOAL",
        font_size=14,
        color=ACCENT_BLUE,
        bold=True,
    )
    add_text_box(
        slide,
        Inches(1.3),
        Inches(2.0),
        Inches(10),
        Inches(0.7),
        "One end-to-end slice: Simulation \u2192 ORCE \u2192 Kafka \u2192 Consumer/Validation",
        font_size=20,
        color=WHITE,
        bold=True,
    )

    # In scope
    add_text_box(
        slide,
        Inches(1),
        Inches(3.2),
        Inches(5),
        Inches(0.5),
        "IN SCOPE",
        font_size=16,
        color=GREEN,
        bold=True,
    )
    add_bullet_list(
        slide,
        Inches(1),
        Inches(3.7),
        Inches(5),
        Inches(3),
        [
            "Structured JSON messages with schema versioning",
            "ISO 8601 UTC timestamps, aligned across feeds",
            "Stable correlation keys (site_id, zone_id, asset_id)",
            "Reproducible deterministic simulation (fixed seed)",
            "TLS-authenticated Kafka transport",
            "9 Kafka topics (5 Smart Energy + 4 Smart City)",
        ],
        font_size=14,
        color=LIGHT_GRAY,
    )

    # Out of scope
    add_text_box(
        slide,
        Inches(7),
        Inches(3.2),
        Inches(5),
        Inches(0.5),
        "OUT OF SCOPE (MS3+)",
        font_size=16,
        color=ORANGE,
        bold=True,
    )
    add_bullet_list(
        slide,
        Inches(7),
        Inches(3.7),
        Inches(5),
        Inches(3),
        [
            "Device control / actuation",
            "Bidirectional communication",
            "Billing / settlement",
            "Real-time safety guarantees",
            "Complex business logic",
        ],
        font_size=14,
        color=LIGHT_GRAY,
    )


def create_architecture_slide(prs):
    """Slide 4: Architecture diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_text_box(
        slide,
        Inches(1),
        Inches(0.5),
        Inches(10),
        Inches(0.8),
        "Architecture: End-to-End Data Flow",
        font_size=36,
        color=WHITE,
        bold=True,
    )
    add_shape(slide, Inches(1), Inches(1.2), Inches(1), Pt(3), ACCENT_BLUE)

    # Architecture boxes
    y_center = Inches(3.5)
    box_h = Inches(2.2)
    arrow_y = y_center + Inches(0.3)

    # Simulation Service box
    sim_box = add_shape(
        slide, Inches(0.5), Inches(2.4), Inches(3), box_h, RGBColor(0x12, 0x2A, 0x4C)
    )
    add_text_box(
        slide,
        Inches(0.7),
        Inches(2.5),
        Inches(2.6),
        Inches(0.4),
        "Simulation Service",
        font_size=18,
        color=ACCENT_BLUE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        Inches(0.7),
        Inches(2.9),
        Inches(2.6),
        Inches(1.6),
        "Python / FastAPI\n\nSmart Energy feeds:\n  Meter, PV, Weather, Price, Consumer\n\nSmart City feeds:\n  Streetlight, Traffic, Event, Visibility\n\nDeterministic RNG (seed=12345)",
        font_size=11,
        color=LIGHT_GRAY,
        alignment=PP_ALIGN.CENTER,
    )

    # Arrow 1
    add_text_box(
        slide,
        Inches(3.5),
        arrow_y,
        Inches(1.2),
        Inches(0.5),
        "HTTP POST\n/api/sim/tick",
        font_size=11,
        color=ORANGE,
        alignment=PP_ALIGN.CENTER,
    )
    add_shape(slide, Inches(3.6), arrow_y + Inches(0.6), Inches(1), Pt(3), ORANGE)

    # ORCE box
    orce_box = add_shape(
        slide, Inches(4.8), Inches(2.4), Inches(3), box_h, RGBColor(0x12, 0x2A, 0x4C)
    )
    add_text_box(
        slide,
        Inches(5.0),
        Inches(2.5),
        Inches(2.6),
        Inches(0.4),
        "ORCE",
        font_size=18,
        color=ACCENT_BLUE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        Inches(5.0),
        Inches(2.9),
        Inches(2.6),
        Inches(1.6),
        "Node-RED (Eclipse XFSC)\n\nSchema validation\nFeed splitting (9 topics)\nKafka publishing\nFlow monitoring\n\nVisual orchestration editor",
        font_size=11,
        color=LIGHT_GRAY,
        alignment=PP_ALIGN.CENTER,
    )

    # Arrow 2
    add_text_box(
        slide,
        Inches(7.8),
        arrow_y,
        Inches(1.2),
        Inches(0.5),
        "Kafka TLS\n9 topics",
        font_size=11,
        color=ORANGE,
        alignment=PP_ALIGN.CENTER,
    )
    add_shape(slide, Inches(7.9), arrow_y + Inches(0.6), Inches(1), Pt(3), ORANGE)

    # Kafka box
    kafka_box = add_shape(
        slide, Inches(9.1), Inches(2.4), Inches(3.5), box_h, RGBColor(0x12, 0x2A, 0x4C)
    )
    add_text_box(
        slide,
        Inches(9.3),
        Inches(2.5),
        Inches(3.1),
        Inches(0.4),
        "Kafka Cluster",
        font_size=18,
        color=ACCENT_BLUE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        Inches(9.3),
        Inches(2.9),
        Inches(3.1),
        Inches(1.6),
        "212.132.83.222:9093\nTLS Client Authentication\n\nsim.smart_energy.*\n  meter, pv, weather, price, consumer\n\nsim.smart_city.*\n  light, traffic, event, weather",
        font_size=11,
        color=LIGHT_GRAY,
        alignment=PP_ALIGN.CENTER,
    )

    # Key properties row
    props_y = Inches(5.2)
    props = [
        ("Deterministic", "Same seed = same output.\nFully reproducible runs.", GREEN),
        ("Correlated", "All feeds share timestamps.\nWeather drives PV output.", LIGHT_BLUE),
        ("Schema-Versioned", "type + schema_version\nin every message.", ORANGE),
        ("TLS-Secured", "Client certificate auth\nfor Kafka transport.", RED_ACCENT),
    ]

    for i, (title, desc, color) in enumerate(props):
        x = Inches(0.5 + i * 3.2)
        add_shape(slide, x, props_y, Inches(2.8), Inches(1.2), RGBColor(0x12, 0x2A, 0x4C))
        add_text_box(
            slide,
            x + Inches(0.15),
            props_y + Inches(0.1),
            Inches(2.5),
            Inches(0.4),
            title,
            font_size=14,
            color=color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_text_box(
            slide,
            x + Inches(0.15),
            props_y + Inches(0.5),
            Inches(2.5),
            Inches(0.6),
            desc,
            font_size=11,
            color=LIGHT_GRAY,
            alignment=PP_ALIGN.CENTER,
        )


def create_smart_energy_slide(prs):
    """Slide 5: Smart Energy Use Case."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_text_box(
        slide,
        Inches(1),
        Inches(0.5),
        Inches(10),
        Inches(0.8),
        "Smart Energy Use Case",
        font_size=36,
        color=WHITE,
        bold=True,
    )
    add_shape(slide, Inches(1), Inches(1.2), Inches(1), Pt(3), ACCENT_BLUE)

    # Feed cards
    feeds = [
        (
            "Energy Meter",
            "sim.smart_energy.meter",
            "site_id, meter_id, timestamp\nactive_power_kw\nactive_energy_kwh_total\n3-phase readings (L1/L2/L3)",
            "2 Janitza UMG96RM meters\nMonotonic energy counter",
        ),
        (
            "PV Generation",
            "sim.smart_energy.pv",
            "site_id, pv_system_id, timestamp\npv_power_kw\ndaily_energy_kwh\nmodule_temperature_c",
            "2 PV systems (10kWp + 5kWp)\nCorrelated with irradiance",
        ),
        (
            "Weather",
            "sim.smart_energy.weather",
            "site_id, timestamp\ntemperature_c\nsolar_irradiance_w_m2\nGHI, DNI, DHI components",
            "Berlin location (52.52, 13.405)\nDrives PV generation",
        ),
        (
            "Energy Price",
            "sim.smart_energy.price",
            "timestamp\nprice_eur_per_kwh\ntariff_type",
            "5 tariff periods per day\nEvening peak: 0.40 EUR/kWh\nNight: 0.15 EUR/kWh",
        ),
        (
            "Consumer Load",
            "sim.smart_energy.consumer",
            "site_id, device_id, timestamp\ndevice_state (ON/OFF)\ndevice_power_kw\ndevice_type",
            "Industrial oven, HVAC,\nCompressor, Pump\n2-3 ON cycles/day",
        ),
    ]

    for i, (title, topic, fields, notes) in enumerate(feeds):
        x = Inches(0.4 + i * 2.55)
        card = add_shape(
            slide, x, Inches(1.6), Inches(2.35), Inches(5.3), RGBColor(0x12, 0x2A, 0x4C)
        )

        add_text_box(
            slide,
            x + Inches(0.1),
            Inches(1.7),
            Inches(2.15),
            Inches(0.4),
            title,
            font_size=14,
            color=ACCENT_BLUE,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_text_box(
            slide,
            x + Inches(0.1),
            Inches(2.05),
            Inches(2.15),
            Inches(0.3),
            topic,
            font_size=9,
            color=ORANGE,
            alignment=PP_ALIGN.CENTER,
        )

        # Separator
        add_shape(slide, x + Inches(0.3), Inches(2.4), Inches(1.75), Pt(1), MEDIUM_GRAY)

        add_text_box(
            slide,
            x + Inches(0.1),
            Inches(2.5),
            Inches(2.15),
            Inches(0.3),
            "Fields",
            font_size=10,
            color=MEDIUM_GRAY,
            bold=True,
        )
        add_text_box(
            slide,
            x + Inches(0.1),
            Inches(2.75),
            Inches(2.15),
            Inches(1.8),
            fields,
            font_size=10,
            color=LIGHT_GRAY,
        )

        add_shape(slide, x + Inches(0.3), Inches(4.6), Inches(1.75), Pt(1), MEDIUM_GRAY)

        add_text_box(
            slide,
            x + Inches(0.1),
            Inches(4.7),
            Inches(2.15),
            Inches(0.3),
            "Details",
            font_size=10,
            color=MEDIUM_GRAY,
            bold=True,
        )
        add_text_box(
            slide,
            x + Inches(0.1),
            Inches(4.95),
            Inches(2.15),
            Inches(1.5),
            notes,
            font_size=10,
            color=LIGHT_GRAY,
        )


def create_smart_city_slide(prs):
    """Slide 6: Smart City Use Case."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_text_box(
        slide,
        Inches(1),
        Inches(0.5),
        Inches(10),
        Inches(0.8),
        "Smart City Use Case",
        font_size=36,
        color=WHITE,
        bold=True,
    )
    add_shape(slide, Inches(1), Inches(1.2), Inches(1), Pt(3), ACCENT_BLUE)

    feeds = [
        (
            "Streetlight",
            "sim.smart_city.light",
            "city_id, zone_id, light_id\ntimestamp\ndimming_level_pct (0-100)\npower_w",
            "6 lights across 2 zones\nSunrise/sunset schedule\nEvent-reactive dimming:\n  Severity 2: +30%\n  Severity 3: +50%",
        ),
        (
            "Traffic",
            "sim.smart_city.traffic",
            "city_id, zone_id\ntimestamp\ntraffic_index (0-100)",
            "Rush hour peaks (07-09, 17-19)\nWeekend reduction (30%)\nPer-zone readings",
        ),
        (
            "City Event",
            "sim.smart_city.event",
            "city_id, zone_id\ntimestamp\nevent_type\nseverity (1-3)\nactive (bool)",
            "Accident, Emergency, Event\nDeterministic scheduling\nEvent mode: 2-3 events/day\nTriggerable dimming change",
        ),
        (
            "Visibility",
            "sim.smart_city.weather",
            "city_id, timestamp\nfog_index (0-100)\nvisibility (good/medium/poor)\nsunrise_time, sunset_time",
            "Dawn fog peak (05-08)\nClears by midday\nSun times from lat/lon",
        ),
    ]

    for i, (title, topic, fields, notes) in enumerate(feeds):
        x = Inches(0.5 + i * 3.1)
        card = add_shape(
            slide, x, Inches(1.6), Inches(2.9), Inches(5.3), RGBColor(0x12, 0x2A, 0x4C)
        )

        add_text_box(
            slide,
            x + Inches(0.15),
            Inches(1.7),
            Inches(2.6),
            Inches(0.4),
            title,
            font_size=16,
            color=ACCENT_BLUE,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_text_box(
            slide,
            x + Inches(0.15),
            Inches(2.1),
            Inches(2.6),
            Inches(0.3),
            topic,
            font_size=10,
            color=ORANGE,
            alignment=PP_ALIGN.CENTER,
        )

        add_shape(slide, x + Inches(0.4), Inches(2.5), Inches(2.1), Pt(1), MEDIUM_GRAY)

        add_text_box(
            slide,
            x + Inches(0.15),
            Inches(2.6),
            Inches(2.6),
            Inches(0.3),
            "Fields",
            font_size=11,
            color=MEDIUM_GRAY,
            bold=True,
        )
        add_text_box(
            slide,
            x + Inches(0.15),
            Inches(2.9),
            Inches(2.6),
            Inches(1.5),
            fields,
            font_size=11,
            color=LIGHT_GRAY,
        )

        add_shape(slide, x + Inches(0.4), Inches(4.5), Inches(2.1), Pt(1), MEDIUM_GRAY)

        add_text_box(
            slide,
            x + Inches(0.15),
            Inches(4.6),
            Inches(2.6),
            Inches(0.3),
            "Details",
            font_size=11,
            color=MEDIUM_GRAY,
            bold=True,
        )
        add_text_box(
            slide,
            x + Inches(0.15),
            Inches(4.9),
            Inches(2.6),
            Inches(1.8),
            notes,
            font_size=11,
            color=LIGHT_GRAY,
        )


def create_correlation_slide(prs):
    """Slide 7: Correlation proof."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_text_box(
        slide,
        Inches(1),
        Inches(0.5),
        Inches(10),
        Inches(0.8),
        "Correlation Proof: MS2 Value Demonstration",
        font_size=36,
        color=WHITE,
        bold=True,
    )
    add_shape(slide, Inches(1), Inches(1.2), Inches(1), Pt(3), ACCENT_BLUE)

    # Left: Smart Energy correlation
    add_shape(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.2), RGBColor(0x12, 0x2A, 0x4C))
    add_text_box(
        slide,
        Inches(0.7),
        Inches(1.7),
        Inches(5.4),
        Inches(0.5),
        "Smart Energy: Cost Impact",
        font_size=20,
        color=GREEN,
        bold=True,
    )

    add_text_box(
        slide,
        Inches(0.7),
        Inches(2.3),
        Inches(5.4),
        Inches(4.2),
        "Given:  meter, PV, weather, price feeds share site_id\n"
        "          and aligned timestamps\n\n"
        "When:  industrial oven runs during MORNING_PEAK\n"
        "          (07:00-09:00, price = 0.33-0.40 EUR/kWh)\n"
        "          vs MIDDAY (11:00-13:00, price = 0.25 EUR/kWh)\n\n"
        "Then:  cost_impact = net_grid_power x price\n"
        "          Morning: high grid import + high price = HIGH cost\n"
        "          Midday: PV offsets load + low price = LOW cost\n\n"
        "Result: 2.7x cost difference is explainable\n"
        "            from price tariff + PV availability",
        font_size=13,
        color=LIGHT_GRAY,
    )

    # Right: Smart City correlation
    add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2), RGBColor(0x12, 0x2A, 0x4C))
    add_text_box(
        slide,
        Inches(7.0),
        Inches(1.7),
        Inches(5.4),
        Inches(0.5),
        "Smart City: Event \u2192 Dimming",
        font_size=20,
        color=LIGHT_BLUE,
        bold=True,
    )

    add_text_box(
        slide,
        Inches(7.0),
        Inches(2.3),
        Inches(5.4),
        Inches(4.2),
        "Given:  event and lighting feeds share zone_id\n"
        "          streetlights follow sunrise/sunset schedule\n\n"
        "When:  emergency event (severity >= 2) occurs\n"
        "          in zone-001\n\n"
        "Then:  streetlights in zone-001 increase dimming:\n"
        "          Severity 2: +30% boost\n"
        "          Severity 3: +50% boost\n"
        "          Change visible within 1 tick (1 minute)\n\n"
        "Result: dimming change is explainable\n"
        "            and repeatable in deterministic mode",
        font_size=13,
        color=LIGHT_GRAY,
    )


def create_demo_slide(prs):
    """Slide 8: Live demo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_text_box(
        slide,
        Inches(1),
        Inches(0.5),
        Inches(10),
        Inches(0.8),
        "Live Demo",
        font_size=36,
        color=WHITE,
        bold=True,
    )
    add_shape(slide, Inches(1), Inches(1.2), Inches(1), Pt(3), ACCENT_BLUE)

    # Demo steps
    steps = [
        (
            "1",
            "Docker Compose Stack",
            "5 services running: Simulation, MQTT, Kafka, ORCE, Kafka UI",
            ACCENT_BLUE,
        ),
        (
            "2",
            "ORCE Flow Editor",
            "Open http://localhost:1880 \u2014 visual flow: HTTP-in \u2192 Validate \u2192 Count \u2192 Respond",
            LIGHT_BLUE,
        ),
        (
            "3",
            "Kafka UI Topics",
            "Open http://localhost:8090 \u2014 9 topics, messages arriving every 60s",
            GREEN,
        ),
        (
            "4",
            "Simulation API",
            "GET /api/v1/health \u2014 service healthy, simulation running",
            ORANGE,
        ),
        (
            "5",
            "Message Inspection",
            "Browse messages in Kafka UI \u2014 show JSON schema, timestamps, correlation keys",
            LIGHT_BLUE,
        ),
        (
            "6",
            "Validation Script",
            "python scripts/demo_e2e.py \u2014 automated schema + correlation validation",
            GREEN,
        ),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        y = Inches(1.8 + i * 0.85)

        num_shape = add_shape(slide, Inches(1), y, Inches(0.45), Inches(0.45), color)
        num_shape.text_frame.paragraphs[0].text = num
        num_shape.text_frame.paragraphs[0].font.size = Pt(16)
        num_shape.text_frame.paragraphs[0].font.color.rgb = WHITE
        num_shape.text_frame.paragraphs[0].font.bold = True
        num_shape.text_frame.paragraphs[0].font.name = "Calibri"
        num_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_text_box(
            slide,
            Inches(1.7),
            y - Inches(0.05),
            Inches(4),
            Inches(0.4),
            title,
            font_size=18,
            color=WHITE,
            bold=True,
        )
        add_text_box(
            slide,
            Inches(1.7),
            y + Inches(0.3),
            Inches(10),
            Inches(0.4),
            desc,
            font_size=14,
            color=LIGHT_GRAY,
        )


def create_validation_slide(prs):
    """Slide 9: Validation results."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_text_box(
        slide,
        Inches(1),
        Inches(0.5),
        Inches(10),
        Inches(0.8),
        "Validation Results",
        font_size=36,
        color=WHITE,
        bold=True,
    )
    add_shape(slide, Inches(1), Inches(1.2), Inches(1), Pt(3), ACCENT_BLUE)

    # BDD checks table
    checks = [
        ("Schema validity", "All messages contain type, schema_version, timestamp", "PASS", GREEN),
        (
            "Timestamp alignment",
            "All feeds share aligned timestamps (1-min intervals)",
            "PASS",
            GREEN,
        ),
        ("Energy counter monotonic", "active_energy_kwh_total strictly increasing", "PASS", GREEN),
        (
            "Cost impact correlation",
            "High-tariff cost > low-tariff cost (2.7x ratio)",
            "PASS",
            GREEN,
        ),
        (
            "Event \u2192 dimming correlation",
            "Severity >= 2 triggers dimming boost in zone",
            "PASS",
            GREEN,
        ),
        (
            "Topic coverage (9/9)",
            "All 9 Kafka topics receive messages at stable cadence",
            "PASS",
            GREEN,
        ),
        (
            "Deterministic reproducibility",
            "Same seed = identical output across runs",
            "PASS",
            GREEN,
        ),
    ]

    # Header
    add_shape(slide, Inches(0.5), Inches(1.6), Inches(12), Inches(0.5), ACCENT_BLUE)
    add_text_box(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(3.5),
        Inches(0.5),
        "Check",
        font_size=14,
        color=WHITE,
        bold=True,
    )
    add_text_box(
        slide,
        Inches(4.2),
        Inches(1.6),
        Inches(6),
        Inches(0.5),
        "Criteria",
        font_size=14,
        color=WHITE,
        bold=True,
    )
    add_text_box(
        slide,
        Inches(10.5),
        Inches(1.6),
        Inches(1.5),
        Inches(0.5),
        "Result",
        font_size=14,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    for i, (check, criteria, result, color) in enumerate(checks):
        y = Inches(2.2 + i * 0.6)
        bg_color = RGBColor(0x12, 0x2A, 0x4C) if i % 2 == 0 else RGBColor(0x0E, 0x22, 0x3E)
        add_shape(slide, Inches(0.5), y, Inches(12), Inches(0.55), bg_color)

        add_text_box(
            slide,
            Inches(0.7),
            y + Inches(0.05),
            Inches(3.3),
            Inches(0.45),
            check,
            font_size=13,
            color=WHITE,
            bold=True,
        )
        add_text_box(
            slide,
            Inches(4.2),
            y + Inches(0.05),
            Inches(6),
            Inches(0.45),
            criteria,
            font_size=12,
            color=LIGHT_GRAY,
        )

        result_box = add_shape(
            slide, Inches(10.7), y + Inches(0.08), Inches(1), Inches(0.35), color
        )
        result_box.text_frame.paragraphs[0].text = result
        result_box.text_frame.paragraphs[0].font.size = Pt(12)
        result_box.text_frame.paragraphs[0].font.color.rgb = WHITE
        result_box.text_frame.paragraphs[0].font.bold = True
        result_box.text_frame.paragraphs[0].font.name = "Calibri"
        result_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Summary box
    add_shape(slide, Inches(0.5), Inches(6.3), Inches(12), Inches(0.7), RGBColor(0x1B, 0x5E, 0x20))
    add_text_box(
        slide,
        Inches(0.7),
        Inches(6.35),
        Inches(11),
        Inches(0.6),
        "MS2 Outcome: Simulation is realistic enough for demo. Data product feasibility is proven.",
        font_size=16,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )


def create_kafka_topics_slide(prs):
    """Slide 10: Kafka topics overview with live counts."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_text_box(
        slide,
        Inches(1),
        Inches(0.5),
        Inches(10),
        Inches(0.8),
        "Kafka Topics: Live Message Counts",
        font_size=36,
        color=WHITE,
        bold=True,
    )
    add_shape(slide, Inches(1), Inches(1.2), Inches(1), Pt(3), ACCENT_BLUE)

    topics = [
        ("sim.smart_energy.meter", "2 meters x N ticks", "meter_id"),
        ("sim.smart_energy.pv", "2 PV systems x N ticks", "pv_system_id"),
        ("sim.smart_energy.weather", "1 station x N ticks", "site_id"),
        ("sim.smart_energy.price", "1 feed x N ticks", "price"),
        ("sim.smart_energy.consumer", "4 devices x N ticks", "device_id"),
        ("sim.smart_city.light", "6 lights x N ticks", "light_id"),
        ("sim.smart_city.traffic", "2 zones x N ticks", "zone_id"),
        ("sim.smart_city.event", "2 zones x N ticks", "zone_id"),
        ("sim.smart_city.weather", "1 city x N ticks", "city_id"),
    ]

    # Header
    add_shape(slide, Inches(1), Inches(1.6), Inches(11), Inches(0.5), ACCENT_BLUE)
    add_text_box(
        slide,
        Inches(1.2),
        Inches(1.6),
        Inches(4.5),
        Inches(0.5),
        "Topic",
        font_size=14,
        color=WHITE,
        bold=True,
    )
    add_text_box(
        slide,
        Inches(5.7),
        Inches(1.6),
        Inches(3),
        Inches(0.5),
        "Message Rate",
        font_size=14,
        color=WHITE,
        bold=True,
    )
    add_text_box(
        slide,
        Inches(8.7),
        Inches(1.6),
        Inches(3),
        Inches(0.5),
        "Partition Key",
        font_size=14,
        color=WHITE,
        bold=True,
    )

    for i, (topic, rate, key) in enumerate(topics):
        y = Inches(2.2 + i * 0.55)
        bg_color = RGBColor(0x12, 0x2A, 0x4C) if i % 2 == 0 else RGBColor(0x0E, 0x22, 0x3E)
        add_shape(slide, Inches(1), y, Inches(11), Inches(0.5), bg_color)

        topic_color = LIGHT_BLUE if "smart_energy" in topic else GREEN
        add_text_box(
            slide,
            Inches(1.2),
            y + Inches(0.05),
            Inches(4.3),
            Inches(0.4),
            topic,
            font_size=13,
            color=topic_color,
            bold=True,
        )
        add_text_box(
            slide,
            Inches(5.7),
            y + Inches(0.05),
            Inches(3),
            Inches(0.4),
            rate,
            font_size=12,
            color=LIGHT_GRAY,
        )
        add_text_box(
            slide,
            Inches(8.7),
            y + Inches(0.05),
            Inches(3),
            Inches(0.4),
            key,
            font_size=12,
            color=MEDIUM_GRAY,
        )


def create_next_steps_slide(prs):
    """Slide 11: Next steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_text_box(
        slide,
        Inches(1),
        Inches(0.5),
        Inches(10),
        Inches(0.8),
        "Next Steps",
        font_size=36,
        color=WHITE,
        bold=True,
    )
    add_shape(slide, Inches(1), Inches(1.2), Inches(1), Pt(3), ACCENT_BLUE)

    # MS2 Completed
    add_shape(slide, Inches(0.5), Inches(1.6), Inches(3.8), Inches(5.2), RGBColor(0x12, 0x2A, 0x4C))
    add_text_box(
        slide,
        Inches(0.7),
        Inches(1.7),
        Inches(3.4),
        Inches(0.5),
        "MS2 (Completed)",
        font_size=18,
        color=GREEN,
        bold=True,
    )
    add_bullet_list(
        slide,
        Inches(0.7),
        Inches(2.3),
        Inches(3.4),
        Inches(4),
        [
            "Simulation service with 9 feeds",
            "ORCE orchestration flows",
            "Kafka TLS transport",
            "Schema validation",
            "Correlation proof (2 domains)",
            "BDD specification (12 scenarios)",
            "Demo script + validation",
        ],
        font_size=13,
        color=LIGHT_GRAY,
    )

    # MS3 Planned
    add_shape(slide, Inches(4.7), Inches(1.6), Inches(3.8), Inches(5.2), RGBColor(0x12, 0x2A, 0x4C))
    add_text_box(
        slide,
        Inches(4.9),
        Inches(1.7),
        Inches(3.4),
        Inches(0.5),
        "MS3 (Planned)",
        font_size=18,
        color=ORANGE,
        bold=True,
    )
    add_bullet_list(
        slide,
        Inches(4.9),
        Inches(2.3),
        Inches(3.4),
        Inches(4),
        [
            "Kafka consumer analytics layer",
            "Dashboard / visualization",
            "Data product API",
            "Advanced correlation analysis",
            "Time-series storage",
            "Performance benchmarking",
        ],
        font_size=13,
        color=LIGHT_GRAY,
    )

    # Future
    add_shape(slide, Inches(8.9), Inches(1.6), Inches(3.8), Inches(5.2), RGBColor(0x12, 0x2A, 0x4C))
    add_text_box(
        slide,
        Inches(9.1),
        Inches(1.7),
        Inches(3.4),
        Inches(0.5),
        "Future",
        font_size=18,
        color=LIGHT_BLUE,
        bold=True,
    )
    add_bullet_list(
        slide,
        Inches(9.1),
        Inches(2.3),
        Inches(3.4),
        Inches(4),
        [
            "Bidirectional control",
            "Real device integration",
            "AI/ML anomaly detection",
            "Multi-site federation",
            "Gaia-X compliance",
            "Production deployment",
        ],
        font_size=13,
        color=LIGHT_GRAY,
    )


def create_closing_slide(prs):
    """Slide 12: Closing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_shape(slide, Inches(1), Inches(2.8), Inches(1.5), Pt(4), ACCENT_BLUE)

    add_text_box(
        slide,
        Inches(1),
        Inches(3.0),
        Inches(11),
        Inches(1),
        "Thank You",
        font_size=44,
        color=WHITE,
        bold=True,
    )

    add_text_box(
        slide,
        Inches(1),
        Inches(4.2),
        Inches(10),
        Inches(1),
        "FACIS FAP  \u2014  IoT & AI over Trusted Zones\nMS2 End-to-End Demonstrator",
        font_size=20,
        color=LIGHT_BLUE,
    )

    add_text_box(
        slide,
        Inches(1),
        Inches(5.5),
        Inches(10),
        Inches(0.8),
        "ATLAS IoT Lab GmbH  |  February 2026",
        font_size=16,
        color=MEDIUM_GRAY,
    )


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    create_title_slide(prs)  # 1
    create_agenda_slide(prs)  # 2
    create_goal_slide(prs)  # 3
    create_architecture_slide(prs)  # 4
    create_smart_energy_slide(prs)  # 5
    create_smart_city_slide(prs)  # 6
    create_correlation_slide(prs)  # 7
    create_demo_slide(prs)  # 8
    create_kafka_topics_slide(prs)  # 9
    create_validation_slide(prs)  # 10
    create_next_steps_slide(prs)  # 11
    create_closing_slide(prs)  # 12

    output = "FACIS_FAP_MS2_Demo_Presentation.pptx"
    prs.save(output)
    print(f"Presentation saved: {output}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
